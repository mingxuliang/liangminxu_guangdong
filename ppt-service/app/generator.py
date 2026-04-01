"""
基于 python-pptx：使用上传的 .pptx 模板中的版式，按大纲逐页生成课件。
模板需至少包含版式：0 通常为标题页，1 通常为「标题+内容」（可含占位符）。
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


def _delete_all_slides(prs: Presentation) -> None:
    """清空所有幻灯片，保留母版与版式。"""
    sld_id_lst = prs.slides._sldIdLst
    while len(sld_id_lst) > 0:
        rId = sld_id_lst[0].rId
        prs.part.drop_rel(rId)
        del sld_id_lst[0]


def _pick_layout(prs: Presentation, index: int) -> Any:
    layouts = prs.slide_layouts
    if index < 0 or index >= len(layouts):
        return layouts[0]
    return layouts[index]


def _set_title_shape(slide: Any, text: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        return
    for shape in slide.shapes:
        if shape.is_placeholder and getattr(shape, "placeholder_format", None):
            try:
                if shape.placeholder_format.idx == 0:
                    shape.text = text
                    return
            except Exception:
                continue


def _fill_body_from_outline(
    slide: Any,
    bullets: list[str],
    sub_items: list[dict[str, Any]] | None,
) -> None:
    """使用「标题+内容」版式时填充正文占位符。"""
    body_shape = None
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        try:
            pf = shape.placeholder_format
            # 1 多为正文；部分模板为 2
            if pf.idx == 1 or (body_shape is None and pf.idx not in (0,)):
                body_shape = shape
        except Exception:
            continue

    if body_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                body_shape = shape
                break

    if body_shape is None or not body_shape.has_text_frame:
        return

    tf = body_shape.text_frame
    tf.clear()
    lines: list[tuple[str, int]] = []

    for b in bullets:
        t = (b or "").strip()
        if t:
            lines.append((t, 0))

    if sub_items:
        for it in sub_items:
            if not isinstance(it, dict):
                continue
            label = (it.get("label") or "").strip()
            if label:
                lines.append((label, 0))
            for ln in it.get("lines") or []:
                if isinstance(ln, str) and ln.strip():
                    lines.append((ln.strip(), 1))

    if not lines:
        tf.text = ""
        return

    tf.text = lines[0][0]
    p0 = tf.paragraphs[0]
    p0.level = lines[0][1]
    p0.font.size = Pt(14)

    for text, level in lines[1:]:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(12 if level > 0 else 14)

    tf.word_wrap = True


def _add_picture_right(slide: Any, image_path: Path, max_w_inch: float = 3.2) -> None:
    """在右侧插入配图（不依赖模板占位）。"""
    try:
        from PIL import Image as PILImage

        with PILImage.open(image_path) as im:
            w_px, h_px = im.size
        ratio = min(max_w_inch * 96 / w_px, 4.5 * 96 / h_px) if w_px and h_px else 1
        w_in = min(max_w_inch, w_px * ratio / 96)
        h_in = h_px * (w_in / (w_px / 96 * 96)) if w_px else 3.0
        h_in = min(h_in, 4.5)
    except Exception:
        w_in, h_in = max_w_inch, 3.0

    left = Inches(6.3)
    top = Inches(1.4)
    slide.shapes.add_picture(str(image_path), left, top, width=Inches(w_in))


def build_presentation(
    template_path: Path,
    outline: dict[str, Any],
    asset_paths: dict[str, Path],
    *,
    title_layout_index: int = 0,
    content_layout_index: int = 1,
) -> Presentation:
    """
    outline 结构示例：
    {
      "course_title": "课程名称",
      "slides": [
        {
          "title": "第一节",
          "bullets": ["要点1", "要点2"],
          "sub_items": [{"label": "活动", "lines": ["步骤a"]}],
          "image_asset_id": "uuid 可选"
        }
      ]
    }
    """
    prs = Presentation(str(template_path))

    try:
        _delete_all_slides(prs)
    except Exception as e:
        logger.warning("清空模板幻灯片失败，尝试继续追加: %s", e)

    course_title = (outline.get("course_title") or "课程课件").strip()
    cover_subtitle = (outline.get("cover_subtitle") or "").strip()
    slides_spec: list[dict[str, Any]] = list(outline.get("slides") or [])
    if not slides_spec:
        slides_spec = [{"title": "内容", "bullets": ["（无大纲条目）"]}]

    title_layout = _pick_layout(prs, title_layout_index)
    content_layout = _pick_layout(prs, content_layout_index)
    if len(prs.slide_layouts) <= content_layout_index:
        content_layout = title_layout

    # 封面（仅 course_title + 可选 cover_subtitle）
    cover = prs.slides.add_slide(title_layout)
    _set_title_shape(cover, course_title)
    if cover_subtitle:
        placed = False
        for shape in cover.shapes:
            if not shape.has_text_frame or not shape.is_placeholder:
                continue
            try:
                idx = shape.placeholder_format.idx
                if idx in (1, 2):
                    shape.text_frame.text = cover_subtitle
                    placed = True
                    break
            except Exception:
                continue
        if not placed:
            title_ph = cover.shapes.title
            for shape in cover.shapes:
                if not shape.has_text_frame:
                    continue
                if title_ph is not None and shape == title_ph:
                    continue
                shape.text_frame.text = cover_subtitle
                break

    # 正文页（与大纲条目一一对应）
    for spec in slides_spec:
        title = (spec.get("title") or "未命名").strip()
        bullets = [str(x).strip() for x in (spec.get("bullets") or spec.get("contentLines") or []) if str(x).strip()]
        sub_items = spec.get("sub_items") or spec.get("subItems")
        if sub_items and not isinstance(sub_items, list):
            sub_items = None

        slide = prs.slides.add_slide(content_layout)
        _set_title_shape(slide, title)
        _fill_body_from_outline(slide, bullets, sub_items)

        aid = spec.get("image_asset_id") or spec.get("imageAssetId")
        if aid and str(aid) in asset_paths:
            p = asset_paths[str(aid)]
            if p.is_file():
                try:
                    _add_picture_right(slide, p)
                except Exception as ex:
                    logger.warning("插入图片失败 %s: %s", p, ex)

    return prs


def save_presentation(prs: Presentation, out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
